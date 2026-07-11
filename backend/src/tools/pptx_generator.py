"""
PowerPoint Generation Tool

Provides create_pptx tool for creating PowerPoint presentations.
Based on Anthropic's Agent Skills for pptx document creation.

Architecture:
1. Agent calls create_pptx with structured spec
2. Tool generates presentation using python-pptx in E2B sandbox
3. File saved via save_generated_file() → Supabase Storage + thread_files
4. Returns structured dict with file_id, download_url, summary
"""

import base64
import json
from typing import Annotated, Any

from langchain_core.runnables import RunnableConfig
from langchain_core.tools import tool

from ..runtime_context import thread_id_ctx, user_id_ctx
from ..sandbox import execute_python

# =============================================================================
# Helper Functions
# =============================================================================


def _resolve_thread_id_from_config(config: RunnableConfig | None) -> str | None:
    """Extract thread_id from RunnableConfig.configurable."""
    if not config:
        return None
    try:
        configurable = config.get("configurable") or {}
        if isinstance(configurable, dict):
            thread_id = configurable.get("thread_id")
            if isinstance(thread_id, str) and thread_id:
                return thread_id
    except Exception:
        return None
    return None


def _resolve_user_id_from_config(config: RunnableConfig | None) -> str | None:
    """Extract user_id from RunnableConfig.configurable."""
    if not config:
        return None
    try:
        configurable = config.get("configurable") or {}
        if isinstance(configurable, dict):
            user_id = configurable.get("user_id")
            if isinstance(user_id, str) and user_id and user_id != "default":
                return user_id
    except Exception:
        return None
    return None


def _validate_context(config: RunnableConfig | None) -> tuple[str, str]:
    """
    Validate and extract user_id and thread_id from config.
    Falls back to ContextVar if not in config.

    Returns:
        Tuple of (user_id, thread_id)
    Raises:
        ValueError if user_id or thread_id cannot be resolved
    """
    user_id = _resolve_user_id_from_config(config)
    thread_id = _resolve_thread_id_from_config(config)

    if not user_id:
        user_id = user_id_ctx.get()
    if not thread_id:
        thread_id = thread_id_ctx.get()

    if not user_id or user_id in ("default", "anonymous"):
        raise ValueError("Unable to identify the user. Please sign in again.")
    if not thread_id or thread_id == "default":
        raise ValueError("Unable to identify the thread. Please refresh the page.")

    return user_id, thread_id


def _save_document_file(
    filename: str,
    data: bytes,
    user_id: str,
    thread_id: str,
) -> dict[str, Any]:
    """
    Save generated document to Supabase Storage and associate with thread.
    """
    from ..storage.file_storage import save_generated_file

    file_info = save_generated_file(
        filename=filename,
        data=data,
        user_id=user_id,
        thread_id=thread_id,
        file_type="generated",
    )

    if not file_info:
        raise RuntimeError(f"Failed to save file {filename} to storage")

    return {
        "file_id": file_info.get("file_id"),
        "download_url": file_info.get("download_url"),
        "filename": filename,
        "file_type": "generated",
    }


# =============================================================================
# PowerPoint Tool
# =============================================================================


@tool
def create_pptx(
    spec: Annotated[
        str,
        """
JSON specification for the PowerPoint presentation:
{
    "title": "Presentation title",
    "slides": [
        {"layout": "title", "title": "Cover title", "subtitle": "Subtitle"},
        {"layout": "content", "title": "Content slide", "bullets": ["Point 1", "Point 2", "Point 3"]},
        {"layout": "two_column", "title": "Comparison", "left": ["Left column content"], "right": ["Right column content"]}
    ]
}
    """,
    ],
    output_name: Annotated[str, "Output filename, such as 'presentation.pptx'"],
    config: RunnableConfig = None,  # type: ignore[assignment]
) -> dict[str, Any]:
    """
    Create a PowerPoint presentation.

    Suitable for business reports, training materials, and project proposals.

    Returns a structured object containing file_id, download_url, and summary.
    """
    try:
        user_id, thread_id = _validate_context(config)
    except ValueError as e:
        return {"success": False, "error": str(e)}

    if not output_name.endswith(".pptx"):
        output_name += ".pptx"

    try:
        spec_dict = json.loads(spec) if isinstance(spec, str) else spec
    except json.JSONDecodeError as e:
        return {"success": False, "error": f"Invalid JSON spec: {e}"}

    code = f'''
import json
import base64
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN

spec = {json.dumps(spec_dict)}
output_name = "{output_name}"

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

for slide_spec in spec.get("slides", []):
    layout_name = slide_spec.get("layout", "content")

    if layout_name == "title":
        # Title slide
        layout = prs.slide_layouts[0]
        slide = prs.slides.add_slide(layout)
        if slide.shapes.title:
            slide.shapes.title.text = slide_spec.get("title", "")
        if len(slide.placeholders) > 1:
            slide.placeholders[1].text = slide_spec.get("subtitle", "")

    elif layout_name == "content":
        # Title and content
        layout = prs.slide_layouts[1]
        slide = prs.slides.add_slide(layout)
        if slide.shapes.title:
            slide.shapes.title.text = slide_spec.get("title", "")

        # Add bullets
        bullets = slide_spec.get("bullets", [])
        if bullets and len(slide.placeholders) > 1:
            tf = slide.placeholders[1].text_frame
            for i, bullet in enumerate(bullets):
                if i == 0:
                    tf.text = bullet
                else:
                    p = tf.add_paragraph()
                    p.text = bullet
                    p.level = 0

    elif layout_name == "two_column":
        # Two column layout
        layout = prs.slide_layouts[3]
        slide = prs.slides.add_slide(layout)
        if slide.shapes.title:
            slide.shapes.title.text = slide_spec.get("title", "")

        # Left column
        left_items = slide_spec.get("left", [])
        if left_items and len(slide.placeholders) > 1:
            tf = slide.placeholders[1].text_frame
            for i, item in enumerate(left_items):
                if i == 0:
                    tf.text = item
                else:
                    p = tf.add_paragraph()
                    p.text = item

        # Right column
        right_items = slide_spec.get("right", [])
        if right_items and len(slide.placeholders) > 2:
            tf = slide.placeholders[2].text_frame
            for i, item in enumerate(right_items):
                if i == 0:
                    tf.text = item
                else:
                    p = tf.add_paragraph()
                    p.text = item

    else:
        # Blank slide with just title
        layout = prs.slide_layouts[5]
        slide = prs.slides.add_slide(layout)
        if slide.shapes.title:
            slide.shapes.title.text = slide_spec.get("title", "")

# Save to bytes
import io
buffer = io.BytesIO()
prs.save(buffer)
file_bytes = buffer.getvalue()

print("__FILE_BASE64_START__")
print(base64.b64encode(file_bytes).decode('utf-8'))
print("__FILE_BASE64_END__")
print(f"Generated {{output_name}} with {{len(spec.get('slides', []))}} slides")
'''

    result = execute_python(code=code, timeout=60, user_id=user_id, thread_id=thread_id)

    if not result.get("success"):
        return {
            "success": False,
            "error": result.get("error") or result.get("stderr", "Execution failed"),
        }

    stdout = result.get("stdout", "")
    try:
        start_idx = stdout.index("__FILE_BASE64_START__") + len("__FILE_BASE64_START__")
        end_idx = stdout.index("__FILE_BASE64_END__")
        b64_content = stdout[start_idx:end_idx].strip()
        file_bytes = base64.b64decode(b64_content)
    except (ValueError, Exception) as e:
        return {"success": False, "error": f"Failed to retrieve file content: {e}"}

    try:
        file_info = _save_document_file(output_name, file_bytes, user_id, thread_id)
    except Exception as e:
        return {"success": False, "error": f"Failed to save file: {e}"}

    slides = spec_dict.get("slides", [])

    return {
        "success": True,
        "file_id": file_info["file_id"],
        "download_url": file_info["download_url"],
        "filename": output_name,
        "file_type": "generated",
        "summary": f"Created PowerPoint presentation {output_name} with {len(slides)} slides",
    }


# =============================================================================
# Export
# =============================================================================

PPTX_TOOLS = [create_pptx]
